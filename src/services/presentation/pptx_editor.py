"""Low-level python-pptx wrapper for creating and persisting presentations."""

import logging
import os
import tempfile
from typing import Optional
from uuid import uuid4

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.slide import SlideLayout

from src.infrastructure.storage.presentation_storage import PresentationStorage

log = logging.getLogger(__name__)


class PptxEditor:
    """Creates, mutates, and saves ``python-pptx`` presentations.

    Wraps the lower-level python-pptx API so that higher-level services can
    work with plain Python objects rather than xml internals.
    """

    def __init__(self, presentation_storage: PresentationStorage) -> None:
        """Initialize the PptxEditor with a PresentationStorage instance."""
        self.presentation_storage = presentation_storage

    def create_blank_presentation(self, template_abs_path: str) -> PptxPresentation:
        """Load a template file and remove all existing slides from it.

        The resulting presentation keeps all slide masters / layouts from the
        template but starts with an empty slide list, ready for new slides to
        be appended.

        Args:
            template_abs_path: Absolute path to the .pptx template file.

        Returns:
            A PptxPresentation with the template theme but no slides.
        """
        prs = PptxPresentation(template_abs_path)
        self._clear_slides(prs)
        return prs

    def add_slide(self, prs: PptxPresentation, layout_index: int) -> Slide:
        """Append a new slide using the layout at *layout_index* and return it.

        Layouts are indexed globally across all slide masters.  When the index
        is out of range the first layout is used as a fallback (see
        :meth:`_get_layout_by_index`).
        """
        layout = self._get_layout_by_index(prs, layout_index)
        return prs.slides.add_slide(layout)

    def save(
        self,
        prs: PptxPresentation,
        presentation_name: Optional[str] = None,
    ) -> str:
        """Persist *prs* to storage and return the storage-relative path.

        The presentation is first written to a temporary file, then handed off
        to :attr:`presentation_storage` for durable storage.  The temp file is
        always cleaned up, even if storage raises.

        Args:
            prs: The presentation to save.
            presentation_name: Desired storage name.  A ``.pptx`` extension is
                added automatically when missing.  Defaults to a UUID-based
                name when ``None``.

        Returns:
            Storage-relative path of the saved file.
        """
        if not presentation_name:
            presentation_name = f"{uuid4()}.pptx"
        elif not presentation_name.lower().endswith(".pptx"):
            presentation_name += ".pptx"

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp_path = tmp.name

        try:
            prs.save(tmp_path)
            return self.presentation_storage.save(
                source_path=tmp_path,
                storage_name=presentation_name,
            )
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    @staticmethod
    def find_placeholder(slide: Slide, ph_type: PP_PLACEHOLDER):
        """Return the first placeholder of *ph_type* on *slide*, or ``None``.

        Args:
            slide: The slide whose placeholders are searched.
            ph_type: A ``PP_PLACEHOLDER`` enum member (e.g. ``TITLE``, ``BODY``).
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.type == ph_type:
                return ph
        return None


    @staticmethod
    def _clear_slides(prs: PptxPresentation) -> None:
        """Remove every slide from prs by detaching its relationship entries.

        Operates directly on the XML slide-ID list so that no slide-specific
        parts remain in the package after this call.
        """
        xml_slides = prs.slides._sldIdLst
        while len(xml_slides):
            slide_id = xml_slides[0]
            r_id = slide_id.get(qn("r:id"))
            prs.part.drop_rel(r_id)
            xml_slides.remove(slide_id)

    @staticmethod
    def _get_layout_by_index(prs: PptxPresentation, layout_index: int) -> SlideLayout:
        """Return the slide layout at the given global layout_index.

        Layouts are numbered contiguously across all masters (master 0's
        layouts first, then master 1's, etc.).  When layout_index exceeds
        the total number of layouts, the first layout of the first master is
        returned and a warning is emitted.
        """
        idx = 0
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if idx == layout_index:
                    return layout
                idx += 1

        log.warning("Layout index %d not found, falling back to index 0", layout_index)
        return prs.slide_masters[0].slide_layouts[0]
