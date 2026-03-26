"""Utilities for writing content (text, images) into python-pptx slide objects."""

import logging
import os

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.slide import Slide
from pptx.util import Inches, Pt

from src.infrastructure.storage.image_storage import ImageStorage
from src.services.presentation.pptx_editor import PptxEditor

log = logging.getLogger(__name__)


class ContentWriter:
    """Writes text and image content into the placeholders and shapes of a slide."""

    def __init__(self, image_storage: ImageStorage) -> None:
        """Initialize the ContentWriter with an ImageStorage instance."""
        self.image_storage = image_storage

    def set_title(self, slide: Slide, title: str) -> bool:
        """Write title into the slide's TITLE or CENTER_TITLE placeholder.

        Returns ``True`` if the title was written successfully, ``False`` otherwise.
        """
        ph = (
            PptxEditor.find_placeholder(slide, PP_PLACEHOLDER.TITLE)
            or PptxEditor.find_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
        )
        if not ph:
            log.debug("No title placeholder found on slide")
            return False

        ph.text = title
        return True

    def set_content(self, slide: Slide, text: str) -> bool:
        """Write text into the slide's BODY or OBJECT placeholder.

        Multi-line text is split on newlines; each segment becomes a separate
        paragraph. Returns ``False`` when no suitable placeholder exists.
        """
        ph = (
            PptxEditor.find_placeholder(slide, PP_PLACEHOLDER.BODY)
            or PptxEditor.find_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        )
        if not ph:
            log.debug("No body/object placeholder found on slide")
            return False

        self._write_paragraphs(ph.text_frame, text)
        return True

    def fill_placeholder(self, slide: Slide, placeholder_idx: int, text: str) -> bool:
        """Write text into the placeholder at placeholder_idx.

        Returns ``True`` if the text was written successfully, ``False`` otherwise.
        """
        try:
            ph = slide.placeholders[placeholder_idx]
        except KeyError:
            return False

        if not ph.has_text_frame:
            return False

        self._write_paragraphs(ph.text_frame, text)
        return True


    def add_image(self, slide: Slide, image_abs_path: str) -> bool:
        """Insert an image into the slide from an absolute file-system path.

        Returns ``True`` if the image was inserted successfully, ``False`` otherwise.
        """
        if not os.path.isfile(image_abs_path):
            log.warning("Image file not found: %s", image_abs_path)
            return False

        ph = PptxEditor.find_placeholder(slide, PP_PLACEHOLDER.PICTURE)
        try:
            ph.insert_picture(image_abs_path)
            return True
        except Exception as exc:
            log.warning("Could not insert picture into placeholder: %s", exc)
            return False


    def add_image_from_storage_path(self, slide: Slide, storage_path: str) -> bool:
        """Resolve storage_path via ImageStorage and insert the image."""
        abs_path = self.image_storage.get_absolute_path(storage_path)
        return self.add_image(slide, abs_path)
    

    @staticmethod
    def _write_paragraphs(text_frame, text: str) -> None:
        """Clear text_frame and populate it with the lines of text."""
        paragraphs = text.split("\n")
        text_frame.clear()
        text_frame.text = paragraphs[0]

        for p_text in paragraphs[1:]:
            p = text_frame.add_paragraph()
            p.text = p_text
