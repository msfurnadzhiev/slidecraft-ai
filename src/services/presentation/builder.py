"""High-level service that assembles a complete presentation from content."""

import logging
import os
from typing import List, Optional

from pptx import Presentation as PptxPresentation
from pptx.slide import Slide

from src.infrastructure.storage.template_storage import TemplateStorage
from src.schemas.presentation.presentation import PresentationContent
from src.schemas.template import SlideLayoutResponse

from src.services.presentation.pptx_editor import PptxEditor
from src.services.presentation.content_writer import ContentWriter
from src.services.presentation.layout_selector import (
    build_layout_map,
    select_layout_index,
)

log = logging.getLogger(__name__)


class PresentationBuilderService:
    """Orchestrates slide creation from a PresentationContent object."""

    def __init__(
        self,
        template_storage: TemplateStorage,
        editor: PptxEditor,
        writer: ContentWriter,
    ) -> None:
        self.template_storage = template_storage
        self.editor = editor
        self.writer = writer


    def create_blank_presentation(self, template_abs_path: str) -> PptxPresentation:
        """Load *template_abs_path* and return a blank presentation."""
        return self.editor.create_blank_presentation(template_abs_path)

    def add_slide(self, prs: PptxPresentation, layout_index: int) -> Slide:
        """Append a new slide with the given layout index and return it."""
        return self.editor.add_slide(prs, layout_index)

    def set_title(self, slide: Slide, title: str) -> bool:
        """Write *title* into the slide's title placeholder."""
        return self.writer.set_title(slide, title)

    def fill_placeholder(self, slide: Slide, placeholder_idx: int, text: str) -> bool:
        """Write *text* into the placeholder at *placeholder_idx*."""
        return self.writer.fill_placeholder(slide, placeholder_idx, text)

    def add_image_from_storage_path(self, slide: Slide, storage_path: str) -> bool:
        """Resolve *storage_path* via ImageStorage and insert the image."""
        return self.writer.add_image_from_storage_path(slide, storage_path)

    def save_presentation(
        self,
        prs: PptxPresentation,
        presentation_name: Optional[str] = None,
    ) -> str:
        """Persist *prs* to storage and return the storage-relative path."""
        return self.editor.save(prs, presentation_name)

    def build_from_content(
        self,
        content: PresentationContent,
        template_file_path: str,
        template_layouts: List[SlideLayoutResponse],
        presentation_name: Optional[str] = None,
    ) -> str:
        """Build a presentation from structured content.
        
        This is the main entry point for generating a `.pptx` file. It:
        1. Loads the template,
        2. Builds a layout map for efficient layout selection,
        3. Iterates through slides and constructs each one,
        4. Saves the final presentation to storage.

        Args:
            content: Structured presentation content produced by the agent pipeline.
            template_file_path: Storage-relative path to the `.pptx` template.
            template_layouts: Layout metadata used to map slide types to layout indices.
            presentation_name: Optional output filename. A UUID is used if not provided.

        Returns:
            Storage-relative path to the saved presentation file.
        """
        prs = self._initialize_presentation(template_file_path)
        layout_map = build_layout_map(template_layouts)

        for slide_content in content.slides:
            self._build_single_slide(prs, slide_content, layout_map)

        return self._finalize_presentation(prs, content, presentation_name)

    def _initialize_presentation(self, template_file_path: str):
        """Load the template and prepare a blank presentation instance."""
        template_abs = self.template_storage.get_absolute_path(template_file_path)
        return self.editor.create_blank_presentation(template_abs)

    def _finalize_presentation(
        self,
        prs,
        content: PresentationContent,
        presentation_name: Optional[str],
    ) -> str:
        """Persist the presentation and log summary information.

        Args:
            prs: The in-memory presentation object.
            content: Source content used to build the presentation.
            presentation_name: Optional output filename.

        Returns:
            Storage-relative path to the saved presentation file.
        """
        storage_path = self.editor.save(prs, presentation_name)

        log.info(
            "Presentation %s saved (%d slides)",
            content.document_id,
            len(content.slides),
        )

        return storage_path

    def _build_single_slide(self, prs, slide_content, layout_map) -> None:
        """Construct and populate a single slide.

        Args:
            prs: The presentation object being built.
            slide_content: Structured data for the current slide.
            layout_map: Precomputed mapping of layout roles to indices.
        """
        layout_index = select_layout_index(
            slide_content.slide_type,
            layout_map,
        )

        slide = self.editor.add_slide(prs, layout_index)

        self._apply_title(slide, slide_content.title)
        self._apply_body(slide, slide_content)
        self._apply_images(slide, slide_content)

        log.debug(
            "Slide %d (%s) built with layout %d",
            slide_content.slide_number,
            slide_content.slide_type,
            layout_index,
        )

    def _apply_title(self, slide, title: str) -> None:
        """Set the slide title if provided."""
        if title:
            self.writer.set_title(slide, title)

    def _apply_body(self, slide, slide_content) -> None:
        """Merge and write body text content into the slide."""
        if not slide_content.content:
            return

        body_text = self._merge_text_content(slide_content)
        if body_text:
            self.writer.set_content(slide, body_text)

    def _apply_images(self, slide, slide_content) -> None:
        """Insert the first valid image found in slide content.

        Images are resolved via storage and only inserted if the file exists.
        """
        if not slide_content.images:
            return

        for img in slide_content.images:
            abs_path = self._resolve_image_path(img.image_url)

            if self._try_add_image(slide, abs_path):
                break

    def _merge_text_content(self, slide_content) -> str:
        """Combine multiple text fragments into a single string."""
        return "\n\n".join(
            tc.text for tc in slide_content.content if tc.text
        )

    def _resolve_image_path(self, image_url: str) -> str:
        """Resolve a storage-relative image path to an absolute filesystem path."""
        return self.writer.image_storage.get_absolute_path(image_url)

    def _try_add_image(self, slide, abs_path: str) -> bool:
        """Attempt to insert an image into the slide."""
        if not os.path.isfile(abs_path):
            return False

        return self.writer.add_image(slide, abs_path)
