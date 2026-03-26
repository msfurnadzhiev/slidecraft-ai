"""Parse PowerPoint (.pptx) template files and extract placeholder metadata."""

import os
from collections import Counter
from typing import Dict, List

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from src.schemas.template import (
    LayoutElementCreate,
    SlideLayoutCreate,
    TemplateContent,
)
from src.utils.singleton import SingletonMeta

# Maps python-pptx placeholder types to semantic roles understood by the agent.
_PLACEHOLDER_ROLE: Dict[int, str] = {
    PP_PLACEHOLDER.TITLE: "title",
    PP_PLACEHOLDER.CENTER_TITLE: "title",
    PP_PLACEHOLDER.SUBTITLE: "subtitle",
    PP_PLACEHOLDER.BODY: "content",
    PP_PLACEHOLDER.OBJECT: "content",
    PP_PLACEHOLDER.PICTURE: "image",
    PP_PLACEHOLDER.CHART: "chart",
    PP_PLACEHOLDER.TABLE: "table",
    PP_PLACEHOLDER.MEDIA_CLIP: "media",
    PP_PLACEHOLDER.ORG_CHART: "diagram",
    PP_PLACEHOLDER.DATE: "date",
    PP_PLACEHOLDER.FOOTER: "footer",
    PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
}


class PPTLoader(metaclass=SingletonMeta):
    """Parse a PowerPoint template and extract slide layout placeholder data."""

    def load_template(self, file_path: str) -> TemplateContent:
        """Open a .pptx file and return a TemplateContent ready for persistence.

        Args:
            file_path: Absolute or relative path to the .pptx file.

        Returns:
            TemplateContent with all extracted layout data.
        """
        prs = Presentation(file_path)
        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)

        layouts: List[SlideLayoutCreate] = []
        seen_names: Counter = Counter()
        layout_index = 0

        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                raw_name = layout.name or f"Layout {layout_index}"
                seen_names[raw_name] += 1
                count = seen_names[raw_name]
                unique_name = raw_name if count == 1 else f"{raw_name} ({count})"

                placeholders = self._extract_placeholders(
                    layout=layout,
                    slide_width=slide_width,
                    slide_height=slide_height,
                )

                layouts.append(
                    SlideLayoutCreate(
                        layout_index=layout_index,
                        name=unique_name,
                        elements=placeholders,
                    )
                )
                layout_index += 1

        return TemplateContent(
            name=os.path.splitext(os.path.basename(file_path))[0],
            file_path=file_path,
            layouts=layouts,
        )

    def _extract_placeholders(
        self,
        layout,
        slide_width: int,
        slide_height: int,
    ) -> List[LayoutElementCreate]:
        """Extract only placeholder shapes from a layout as normalised elements.
        
        Args:
            layout: The layout to extract placeholders from.
            slide_width: The width of the slide.
            slide_height: The height of the slide.

        Returns:
            A list of LayoutElementCreate objects.
        """
        elements: List[LayoutElementCreate] = []

        for shape in layout.shapes:
            if not shape.is_placeholder:
                continue

            ph = shape.placeholder_format
            elements.append(
                LayoutElementCreate(
                    placeholder_idx=ph.idx,
                    role=_PLACEHOLDER_ROLE.get(ph.type),
                    x=shape.left / slide_width,
                    y=shape.top / slide_height,
                    width=shape.width / slide_width,
                    height=shape.height / slide_height,
                )
            )

        return elements
